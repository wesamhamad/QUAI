from pptx import Presentation

from .base import BaseExtractor


class PptxExtractor(BaseExtractor):
    def extract(self, file_path: str) -> dict:
        prs = Presentation(file_path)
        slides_text = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_parts = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_parts.append(text)

                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text.strip()
                            for cell in row.cells
                            if cell.text.strip()
                        )
                        if row_text:
                            slide_parts.append(row_text)

            if slide_parts:
                header = f"--- Slide {slide_num} / \u0634\u0631\u064a\u062d\u0629 {slide_num} ---"
                slides_text.append(f"\n{header}\n" + "\n".join(slide_parts))

        full_text = "\n".join(slides_text)

        return {
            "content": full_text,
            "page_count": len(prs.slides),
            "ocr_applied": False,
            "language_detected": self.detect_language(full_text),
        }
